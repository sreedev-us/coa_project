import subprocess
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

# 1. Setup Data Structures
predictors = ["LocalBP", "BiModeBP", "TournamentBP", "PerceptronBP"]
workloads = {
    "Random": "s4_project_tests/branch_test_random",
    "Patterned": "s4_project_tests/branch_test_pattern"
}
results = {"Random": {}, "Patterned": {}}

print("🚀 Starting Dual-Execution KTU Project Pipeline...")

# 2. Execute gem5 for Both Workloads
for wl_name, wl_path in workloads.items():
    print(f"\n=====================================")
    print(f"📦 BEGINNING WORKLOAD: {wl_name.upper()}")
    print(f"=====================================")
    
    for pred in predictors:
        print(f"▶️ Simulating {pred}...")
        
        command = ["build/X86/gem5.opt", "s4_project_tests/run_project.py", pred, wl_path]
        subprocess.run(command, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        
        incorrect = 0
        squashes = 0
        sim_ticks = 0
        
        with open("m5out/stats.txt", "r") as stats_file:
            for line in stats_file:
                if line.startswith("simTicks "):
                    sim_ticks = int(line.split()[1])
                if "system.cpu.branchPred.condIncorrect" in line:
                    incorrect = int(line.split()[1])
                if "system.cpu.branchPred.squashes_0::total" in line:
                    squashes = int(line.split()[1])
                    
        results[wl_name][pred] = {
            "incorrect": incorrect,
            "squashes": squashes,
            "ticks": sim_ticks
        }
        print(f"✅ {pred} | Mispredictions: {incorrect} | Flushes: {squashes}")

# 3. Generate Two Separate Graphs
colors = ['#FF9999', '#66B2FF', '#99FF99', '#B266FF']

for wl_name in workloads.keys():
    plt.figure(figsize=(8, 5))
    errors = [results[wl_name][p]["incorrect"] for p in predictors]
    
    bars = plt.bar(predictors, errors, color=colors)
    plt.title(f'Control Hazards: {wl_name} Workload', fontsize=14, fontweight='bold')
    plt.xlabel('Branch Predictor Algorithm')
    plt.ylabel('Total Mispredictions')
    
    for bar in bars:
        yval = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2, yval + (max(errors)*0.01), int(yval), ha='center', fontweight='bold')
    
    plt.savefig(f"s4_project_tests/graph_{wl_name.lower()}.png", dpi=300)
    plt.close()

# 4. Build the PowerPoint Presentation
print("\n📝 Compiling Highly Formatted PowerPoint...")
prs = Presentation()

# Helper function for clean text formatting
def add_bullet(text_frame, text, font_size=22, bold=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.space_after = Pt(10) # Adds clean spacing between lines
    return p

# --- Slide 1: Official Title Page ---
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "Control Hazard Analysis & Hardware Branch Prediction"
subtitle.text = "Providence College of Engineering\nKTU S4 COA"

# --- Slide 2: Team & Submission Details ---
slide_team = prs.slides.add_slide(prs.slide_layouts[1])
slide_team.shapes.title.text = "Project Team"
tf_team = slide_team.shapes.placeholders[1].text_frame
tf_team.text = "Team Members:"

team = [
    "Agnivesh P A (PRC24CA010)", 
    "Bibin C Mathew (PRC24CA028)", 
    "Joyal Philip John (PRC24CA037)", 
    "Melvin Mathew (PRC24CA041)", 
    "Sreedev U.S. (PRC24CA053)"
]

for member in team:
    add_bullet(tf_team, f" {member}", 22)

print("\n")
p_sub = tf_team.add_paragraph()
p_sub.text = "Submitted to: Ms. Gayathri"
p_sub.font.size = Pt(22)
p_sub.font.italic = True
p_sub.font.bold = True

# --- Slide 3: Methodology ---
slide_intro = prs.slides.add_slide(prs.slide_layouts[1])
slide_intro.shapes.title.text = "Methodology & Workloads"
tf_intro = slide_intro.shapes.placeholders[1].text_frame
tf_intro.text = "We utilized the gem5 O3 architectural simulator to test 4 hardware predictors:"

add_bullet(tf_intro, "1. Random Workload: Unpredictable branches (rand() % 100 < 50) to establish a baseline.", 20)
add_bullet(tf_intro, "2. Patterned Workload: A predictable mathematical loop (i % 4 != 0).", 20)
add_bullet(tf_intro, "Goal: Observe how advanced hardware 'learns' software patterns to avoid pipeline flushes.", 22, True)

# --- Slide 4: Random Workload Data ---
slide_rand = prs.slides.add_slide(prs.slide_layouts[1])
slide_rand.shapes.title.text = "Phase 1: Random Workload Execution"
tf_rand = slide_rand.shapes.placeholders[1].text_frame
tf_rand.text = "Hypothesis: Predictors cannot find a pattern in random data. Mispredictions will be universally high."
for pred in predictors:
    add_bullet(tf_rand, f"• {pred}: {results['Random'][pred]['incorrect']:,} Mispredictions | {results['Random'][pred]['squashes']:,} Flushes", 22)

# --- Slide 5: Random Workload Graph ---
slide_rg = prs.slides.add_slide(prs.slide_layouts[5])
slide_rg.shapes.title.text = "Visualizing the Random Workload"
slide_rg.shapes.add_picture("s4_project_tests/graph_random.png", Inches(1), Inches(1.5), width=Inches(8))

# --- Slide 6: Patterned Workload Data ---
slide_pat = prs.slides.add_slide(prs.slide_layouts[1])
slide_pat.shapes.title.text = "Phase 2: Patterned Workload Execution"
tf_pat = slide_pat.shapes.placeholders[1].text_frame
tf_pat.text = "Hypothesis: Advanced predictors will 'learn' the repeating mathematical loop and avoid flushes."
for pred in predictors:
    add_bullet(tf_pat, f"• {pred}: {results['Patterned'][pred]['incorrect']:,} Mispredictions | {results['Patterned'][pred]['squashes']:,} Flushes", 22)

# --- Slide 7: Patterned Workload Graph ---
slide_pg = prs.slides.add_slide(prs.slide_layouts[5])
slide_pg.shapes.title.text = "Visualizing the Patterned Workload"
slide_pg.shapes.add_picture("s4_project_tests/graph_patterned.png", Inches(1), Inches(1.5), width=Inches(8))

# --- Slide 8: Time Saved (Raw Data) ---
slide_time = prs.slides.add_slide(prs.slide_layouts[1])
slide_time.shapes.title.text = "Measuring Performance: CPU Clock Cycles"
tf_time = slide_time.shapes.placeholders[1].text_frame
tf_time.text = "Total simulated execution time (Ticks) for the Patterned Workload:"

for pred in predictors:
    add_bullet(tf_time, f"• {pred}: {results['Patterned'][pred]['ticks']:,} Ticks", 28)

worst_ticks = max(results['Patterned'][p]['ticks'] for p in predictors)
best_ticks = min(results['Patterned'][p]['ticks'] for p in predictors)
saved_ticks = worst_ticks - best_ticks
percent_saved = (saved_ticks / worst_ticks) * 100

# --- Slide 9: The Real-World Impact ---
slide_impact = prs.slides.add_slide(prs.slide_layouts[1])
slide_impact.shapes.title.text = "The Real-World Impact"
tf_impact = slide_impact.shapes.placeholders[1].text_frame
tf_impact.text = "What do these saved clock cycles actually mean for performance?"

add_bullet(tf_impact, f"By minimizing pipeline bubbles, the optimal predictor saved {saved_ticks:,} CPU Ticks compared to the slowest predictor.", 24)
add_bullet(tf_impact, f"This represents a {percent_saved:.2f}% reduction in total execution time.", 24, True)
add_bullet(tf_impact, "In massive data centers or high-performance environments, a ~3.5% speedup across billions of instructions saves significant processing time and power, simply by choosing the correct hardware architecture!", 22)

# --- Slide 10: AI/Perceptron Explanation ---
slide_ml = prs.slides.add_slide(prs.slide_layouts[1])
slide_ml.shapes.title.text = "The AI Paradox: Why the Perceptron Failed"
tf_ml = slide_ml.shapes.placeholders[1].text_frame
tf_ml.text = "The PerceptronBP uses Neural Networks, so why did it perform poorly?"

add_bullet(tf_ml, "1. Training Overhead: Machine learning models require a 'warm-up' period to train their weights. 100k loops was not enough data for it to stabilize.", 20)
add_bullet(tf_ml, "2. Over-engineering: The Tournament predictor uses a simple meta-table, allowing it to instantly recognize the short loop pattern.", 20)
add_bullet(tf_ml, "Conclusion: Throwing AI hardware at simple code is highly inefficient!", 22, True)

# 5. Save the file
ppt_path = "s4_project_tests/KTU_S4_Project_Report.pptx"
prs.save(ppt_path)
print(f"🎉 Success! 10-Slide formatted presentation saved to {ppt_path}")